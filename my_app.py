import streamlit as st
import google.generativeai as genai
import pdfplumber
import os
import re
import time
import json
from datetime import date
import io # Used to save files in memory
from docx import Document # For Word docs
from docx.shared import Pt # For setting font sizes
from docx.enum.text import WD_ALIGN_PARAGRAPH # For alignment

# --- (1) PERSISTENT FILE HELPERS ---

def load_from_file(filename, default_data):
    """Loads data from a JSON file. If file doesn't exist or is old, creates a new one."""
    today_str = str(date.today())
    
    if not os.path.exists(filename):
        save_to_file(filename, default_data)
        return default_data
    
    try:
        with open(filename, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        # Auto-reset daily USAGE counter
        if filename == "usage.json" and data.get("date") != today_str:
            # We can't use st.success() here, as it might run before st.set_page_config
            print("New day! Resetting API counters.") 
            save_to_file(filename, default_data)
            return default_data
        
        # Migrate old usage file format (if necessary)
        if filename == "usage.json" and "count" in data and "counts" not in data:
            new_data = {
                "date": data.get("date", today_str),
                "counts": {"total_legacy_calls": data.get("count", 0)}
            }
            save_to_file(filename, new_data)
            return new_data

        return data
            
    except json.JSONDecodeError:
        # We can't use st.error() here
        print(f"Error reading {filename}. File might be corrupt. Creating a new one.")
        save_to_file(filename, default_data)
        return default_data

def save_to_file(filename, data):
    """Saves data to a JSON file."""
    with open(filename, 'w', encoding='utf-8') as f:
        json.dump(data, f, indent=4, ensure_ascii=False)

# --- (2) CORE GEMINI API FUNCTION ---

DELAYS = {
    "gemini-2.5-flash-lite": 5,
    "gemini-2.5-flash": 7,
    "gemini-2.5-pro": 13
}

def call_gemini_api(system_prompt, user_prompt, temperature, model_name, chat_history=None):
    """
    A single, safe function to call the Gemini API.
    It checks for an API key, updates the usage counter, and respects rate limits.
    """
    api_key = st.session_state.get("api_key")
    if not api_key:
        st.error("API Key not set. Please enter your API key in the sidebar.")
        return None

    try:
        genai.configure(api_key=api_key)
        config = genai.GenerationConfig(temperature=temperature)
        model = genai.GenerativeModel(
            model_name=model_name,
            system_instruction=system_prompt,
            generation_config=config
        )
        
        if chat_history:
            chat = model.start_chat(history=chat_history)
            response = chat.send_message(user_prompt)
        else:
            response = model.generate_content(user_prompt)
        
        if model_name not in st.session_state.usage_data["counts"]:
            st.session_state.usage_data["counts"][model_name] = 0
        st.session_state.usage_data["counts"][model_name] += 1
        save_to_file("usage.json", st.session_state.usage_data)
        
        delay = DELAYS.get(model_name, 7)
        time.sleep(delay)
        
        return response.text

    except Exception as e:
        st.error(f"An API error occurred: {e}")
        return None

def create_toc_from_headings(html_content):
    """Scans final HTML for h2 tags and builds a ToC."""
    st.write("    > Building Table of Contents...")
    headings = re.findall(
        r'<h2\s+id=["\'](.*?)["\']>(.*?)</h2>', 
        html_content,
        re.IGNORECASE
    )
    if not headings:
        st.write("    > No h2 headings with IDs found. Skipping ToC.")
        return ""
    toc_lines = ['<details class="toc" open><summary>Table of Contents</summary><ul>']
    for id, text in headings:
        toc_lines.append(f'<li><a href="#{id}">{text}</a></li>')
    toc_lines.append('</ul></details><br><hr><br>')
    st.write(f"    > ToC created with {len(headings)} entries.")
    return "\n".join(toc_lines)

# --- (3) DEFAULT DATA (FOR NEW FILES) ---
DEFAULT_GEMS = {
    "Blank Chat Prompt": "You are a helpful assistant.",
    
    "Test Generator (.docx)": """
You are an expert VCE Science exam designer.
Your task is to generate a list of Multiple Choice Questions (MCQs) and Short Answer Questions (SAQs) based on a topic and a rubric.
Your output MUST be a single, valid JSON object. Do not include ```json backticks or any other text.

**CRITICAL INSTRUCTIONS:**
1.  **Rubric Coverage:** You *must* generate at least one question that assesses each and every criterion in the provided rubric.
2.  **Question Types:** Generate the exact number of MCQs and SAQs requested.
3.  **Strict JSON:** The entire output must be a single JSON object.

**JSON FORMAT:**
{
  "mcqs": [
    {
      "question_text": "The tendency to attribute our successes to internal factors and failures to external factors is called:",
      "options": [
        "A. Fundamental attribution error",
        "B. Actor-observer bias",
        "C. Self-serving bias",
        "D. Cognitive dissonance"
      ]
    }
  ],
  "saqs": [
    {
      "question_text": "Explain the difference between the 'affective' and 'behavioural' components of an attitude, providing an example for each.",
      "marks": 4
    }
  ]
}
""",
    "PowerPoint Generator (.pptx)": """
You are an expert VCE educator. Your task is to generate the *content* for a PowerPoint presentation based on a textbook chunk.
The output format MUST be a specific JSON object. Do not include ```json backticks.
Your entire response must be *only* the JSON object.
**EXAMPLE:**
{
  "slides": [
    {"title": "Slide 1 Title", "body": ["Bullet point 1", "Bullet point 2"]},
    {"title": "Slide 2 Title", "body": ["Bullet point 1", "Bullet point 2", "Bullet point 3"]}
  ]
}
""",
    "Gimkit Generator (.csv)": """
You are a question generator. Your task is to create a list of questions and answers on a given topic.
The output format MUST be a valid CSV (Comma Separated Values) text.
Do not add any other text, explanation, or ```csv backticks.
Your entire response must be *only* the raw CSV data.
**EXAMPLE OUTPUT:**
"What is the capital of France?","Paris"
"Who wrote Hamlet?","William Shakespeare"
"What is 2+2?","4"
""",
    "Rubric Comment Generator": """
You are an expert VCE teacher, skilled at writing constructive feedback.
You will be given:
1.  The test questions.
2.  The rubric.
3.  A description of what the student did correctly and incorrectly.
Do not add any other text, just the final comment.

**YOUR TASK:**
Generate a concise, constructive comment (1-2 paragraphs) for a student report.
-   Start with what the student did well, referencing the rubric.
-   Clearly explain what they missed or misunderstood, referencing the questions.
-   Provide a clear, actionable "next step" for improvement.
-   Maintain a professional and encouraging tone.
"""
}
DEFAULT_USAGE = {
    "date": str(date.today()),
    "counts": {
        "gemini-2.5-flash-lite": 0,
        "gemini-2.5-flash": 0,
        "gemini-2.5-pro": 0
    }
}
DEFAULT_CHATS = {}

OUTPUT_FOLDER = r'C:\Users\hgh\OneDrive - Brentwood Secondary College\Desktop\Textbook_HTML_Files'

# --- (4) NEW, ROBUST WORD DOC HELPER FUNCTIONS ---

def find_and_replace(doc, find_str, replace_str):
    """Finds and replaces text in all paragraphs and tables in a .docx"""
    for para in doc.paragraphs:
        if find_str in para.text:
            para.text = para.text.replace(find_str, replace_str)
            # Re-apply center alignment if it's a title placeholder
            if find_str in ["{{SUBJECT_TITLE}}", "{{UNIT_TITLE}}", "{{YEAR}}"]:
                para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    if find_str in para.text:
                        para.text = para.text.replace(find_str, replace_str)

def inject_mcq_questions(doc, mcqs):
    """Finds the {{MCQ_SECTION}} placeholder and injects MCQs"""
    for para in doc.paragraphs:
        if "{{MCQ_SECTION}}" in para.text:
            para.text = "" # Clear the placeholder
            
            p_section = doc.add_paragraph()
            p_section.add_run("SECTION A – Multiple-Choice Questions").bold = True
            
            p_instructions = doc.add_paragraph()
            p_instructions.add_run("Answer all questions by selecting the correct option.")
            
            for i, q_data in enumerate(mcqs):
                doc.add_paragraph() # Add a space
                q_para = doc.add_paragraph()
                run = q_para.add_run(f"Question {i+1}\n")
                run.bold = True
                run = q_para.add_run(q_data.get("question_text", ""))
                run.bold = True
                
                for opt in q_data.get("options", []):
                    doc.add_paragraph(opt, style='List Paragraph')
            return True # Success
    return False # Placeholder not found

def inject_saq_questions(doc, saqs):
    """Finds the {{SAQ_SECTION}} placeholder and injects SAQs"""
    for para in doc.paragraphs:
        if "{{SAQ_SECTION}}" in para.text:
            para.text = "" # Clear the placeholder
            
            p_section = doc.add_paragraph()
            p_section.add_run("SECTION B – Short-Answer Questions").bold = True
            
            p_instructions = doc.add_paragraph()
            p_instructions.add_run("Answer all questions in the spaces provided.")
            
            for i, q_data in enumerate(saqs):
                doc.add_paragraph() # Add a space
                q_para = doc.add_paragraph()
                
                marks = q_data.get("marks", 1)
                run = q_para.add_run(f"Question {i+1} ({marks} mark{'s' if marks > 1 else ''})\n")
                run.bold = True
                run = q_para.add_run(q_data.get("question_text", ""))
                run.bold = True
                
                # Add blank lines for answer space
                for _ in range(marks * 3):
                    doc.add_paragraph("_________________________________________________________________")
            
            return True # Success
    return False # Placeholder not found


# --- (5) THIS IS THE MAIN FUNCTION THAT RUNS THE APP ---
def main():
    """Main function to run the Streamlit app."""
    
    # --- (A) PAGE CONFIGURATION ---
    # !! MUST BE THE FIRST STREAMLIT COMMANDS !!
    st.set_page_config(
        page_title="VCE AI Teacher Toolkit",
        layout="wide"
    )

    st.markdown("""
    <style>
        .toc { background-color: #f9f9f9; border: 1px solid #ddd; border-radius: 8px; padding: 15px 25px; margin-bottom: 25px; }
        .toc summary { font-size: 1.2em; font-weight: bold; cursor: pointer; }
        .toc ul { margin-top: 10px; }
        .toc li { margin-bottom: 5px; }
        .toc a { text-decoration: none; color: #0066cc; }
        .toc a:hover { text-decoration: underline; }
    </style>
    """, unsafe_allow_html=True)
    
    # --- (B) SESSION STATE INITIALIZATION ---
    # This can run now, after set_page_config
    
    if "api_key" not in st.session_state:
        st.session_state.api_key = None

    if "gems" not in st.session_state:
        st.session_state.gems = load_from_file("gems.json", DEFAULT_GEMS)
        
        # Migration Check: Add missing default Gems
        migrated = False
        for gem_name, gem_prompt in DEFAULT_GEMS.items():
            if gem_name not in st.session_state.gems:
                st.session_state.gems[gem_name] = gem_prompt
                migrated = True
        
        if migrated:
            save_to_file("gems.json", st.session_state.gems)
            # This st.info() call is now SAFE
            st.info("Added missing default Gems to your 'gems.json' file!") 
            time.sleep(2)
            st.rerun()

    if "usage_data" not in st.session_state:
        st.session_state.usage_data = load_from_file("usage.json", DEFAULT_USAGE)

    if "chats" not in st.session_state:
        st.session_state.chats = load_from_file("chats.json", DEFAULT_CHATS)

    if "current_chat_id" not in st.session_state:
        st.session_state.current_chat_id = None
        
    # --- (C) OTHER SETUP ---
    # This can run now
    os.makedirs(OUTPUT_FOLDER, exist_ok=True)


    # --- (D) UI: SIDEBAR ---
    st.sidebar.title("VCE AI Teacher Toolkit")
    st.sidebar.markdown("Welcome! This app helps you generate VCE resources. **Please enter your own Google API key below to get started.**")

    st.session_state.api_key = st.sidebar.text_input(
        "Enter your Google API Key",
        type="password",
        value=st.session_state.get("api_key")
    )

    model_name = st.sidebar.selectbox(
        "Choose your AI model",
        ("gemini-2.5-flash-lite", "gemini-2.5-flash", "gemini-2.5-pro"),
        key="model_name"
    )

    temperature = st.sidebar.slider(
        "Set AI 'creativity' (Temperature)",
        min_value=0.0, max_value=2.0, value=st.session_state.get("temperature", 0.0),
        step=0.1, key="temperature",
        help="0.0 = Factual, 2.0 = Creative"
    )

    st.sidebar.markdown("---")
    st.sidebar.subheader("Daily Usage")
    counts = st.session_state.usage_data.get("counts", {})
    LIMITS = {"gemini-2.5-flash-lite": 1000, "gemini-2.5-flash": 250, "gemini-2.5-pro": 50}

    if not counts:
        st.sidebar.text("No calls made today.")
    else:
        for model, count in counts.items():
            limit = LIMITS.get(model, 1)
            label = f"{model} ({count} / {limit})"
            progress = min(count / limit, 1.0)
            st.sidebar.text(label)
            st.sidebar.progress(progress)

    if st.sidebar.button("Reset All Counters"):
        st.session_state.usage_data = DEFAULT_USAGE
        save_to_file("usage.json", st.session_state.usage_data)
        st.rerun()

    st.sidebar.markdown("---")
    st.sidebar.subheader("Saved Chats")

    if st.sidebar.button("New Chat", use_container_width=True):
        st.session_state.current_chat_id = None
        st.rerun()

    sorted_chat_ids = sorted(st.session_state.chats.keys(), reverse=True)

    for chat_id in sorted_chat_ids:
        chat_title = st.session_state.chats[chat_id]["title"]
        if st.sidebar.button(chat_title, key=f"chat_{chat_id}", use_container_width=True):
            st.session_state.current_chat_id = chat_id
            st.rerun()

    # --- (E) UI: MAIN PAGE (TABS) ---
    st.title("VCE Resource Generators")
    
    tab_test, tab_ppt, tab_gimkit, tab_comment, tab_gems = st.tabs([
        "Test Generator (.docx)", 
        "PPT Generator (.pptx)", 
        "Gimkit Generator (.csv)", 
        "Comment Generator", 
        "Gem Creator"
    ])

    # --- TAB 1: TEST GENERATOR ---
    with tab_test:
        st.header("Test Generator (.docx)")
        st.info("This tool uses a `test_template.docx` file in your app folder. Please update it with the new placeholders!")
        
        gem_name = "Test Generator (.docx)"
        
        with st.expander("View/Edit Gem Prompt"):
            st.session_state.gems[gem_name] = st.text_area(
                "Prompt for Test Generator:",
                value=st.session_state.gems.get(gem_name),
                height=300,
                key="gem_test"
            )
        
        st.subheader("1. Fill in Title Page Details")
        col1, col2, col3 = st.columns(3)
        with col1:
            subject_title = st.text_input("Subject Title", "Year 7 Science")
        with col2:
            unit_title = st.text_input("Unit/SAC Title", "Unit 2 SAC Outcome 1")
        with col3:
            year = st.text_input("Year", "2025")

        st.subheader("2. Define Test Structure")
        col1, col2 = st.columns(2)
        with col1:
            num_mcq = st.number_input("Number of MCQs", min_value=0, value=10)
            marks_mcq = st.number_input("Total Marks for MCQs", min_value=0, value=10)
        with col2:
            num_saq = st.number_input("Number of SAQs", min_value=0, value=5)
            marks_saq = st.number_input("Total Marks for SAQs", min_value=0, value=10)
        
        st.subheader("3. Provide Content & Rubric")
        topic = st.text_input("Topic for the test:", "Year 7 Biology")
        rubric_text = st.text_area("Paste Rubric Here (to ensure question coverage):", "e.g., 'Criterion 1: Defines key terms.'\n'Criterion 2: Applies concepts to scenarios.'")

        if st.button("Generate Test"):
            if not all([topic, rubric_text, subject_title, unit_title, year]):
                st.warning("Please fill in all fields to generate the test.")
            elif not st.session_state.api_key:
                st.error("API Key not set. Please enter your API key in the sidebar.")
            else:
                user_prompt = f"""
                Topic: {topic}
                Rubric: {rubric_text}
                Number of MCQs: {num_mcq}
                Number of SAQs: {num_saq}
                Total Marks for SAQs: {marks_saq}
                """
                with st.spinner("Calling Gemini to generate questions (as JSON)..."):
                    response = call_gemini_api(
                        st.session_state.gems[gem_name],
                        user_prompt,
                        st.session_state.temperature,
                        st.session_state.model_name
                    )
                
                if response:
                    try:
                        st.write("  > AI returned content. Parsing JSON...")
                        
                        json_start = response.find('{')
                        json_end = response.rfind('}') + 1
                        if json_start == -1 or json_end == -1:
                            raise json.JSONDecodeError("No JSON object found in response.", response, 0)
                        response_clean = response[json_start:json_end]
                        data = json.loads(response_clean)
                        
                        st.write("  > JSON parsed. Opening `test_template.docx`...")
                        
                        template_path = "test_template.docx"
                        if not os.path.exists(template_path):
                            st.error(f"`{template_path}` not found! Please create it and add the placeholders (e.g., {{SUBJECT_TITLE}}).")
                            st.stop()
                            
                        doc = Document(template_path)
                        
                        st.write("  > Replacing placeholders...")
                        total_questions = num_mcq + num_saq
                        total_marks = marks_mcq + marks_saq
                        
                        replacements = {
                            "{{SUBJECT_TITLE}}": subject_title,
                            "{{SUBJECT TITLE}}": subject_title,
                            "{{UNIT_TITLE}}": unit_title,
                            "{{Unit_Title}}": unit_title,
                            "{{YEAR}}": year,
                            "{{MCQ_NUM}}": str(num_mcq),
                            "{{MCQ_MARKS}}": str(marks_mcq),
                            "{{SAQ_NUM}}": str(num_saq),
                            "{{SAQ_MARKS}}": str(marks_saq),
                            "{{TOTAL_QUESTIONS}}": str(total_questions),
                            "{{TOTAL_MARKS}}": str(total_marks),
                        }
                        
                        for para in doc.paragraphs:
                            for key, value in replacements.items():
                                if key in para.text:
                                    para.text = para.text.replace(key, value)
                                    if key in ["{{SUBJECT_TITLE}}", "{{UNIT_TITLE}}", "{{YEAR}}", "{{SUBJECT TITLE}}", "{{Unit_Title}}"]:
                                        para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        
                        for table in doc.tables:
                            for row in table.rows:
                                for cell in row.cells:
                                    for para in cell.paragraphs:
                                        for key, value in replacements.items():
                                            if key in para.text:
                                                para.text = para.text.replace(key, value)

                        st.write("  > Injecting MCQs and SAQs...")
                        mcq_success = inject_mcq_questions(doc, data.get("mcqs", []))
                        saq_success = inject_saq_questions(doc, data.get("saqs", []))
                        
                        if not mcq_success:
                            st.warning("Could not find '{{MCQ_SECTION}}' placeholder. MCQs were not added.")
                        if not saq_success:
                            st.warning("Could not find '{{SAQ_SECTION}}' placeholder. SAQs were not added.")

                        st.write("  > Saving to memory buffer...")
                        bio = io.BytesIO()
                        doc.save(bio)
                        
                        st.success("Test generated!")
                        st.download_button(
                            label="Download Test as .docx",
                            data=bio.getvalue(),
                            file_name=f"{subject_title.replace(' ', '_')}_{unit_title.replace(' ', '_')}_test.docx",
                            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                        )
                        
                    except json.JSONDecodeError:
                        st.error("AI returned invalid JSON. Could not create .docx.")
                        st.text_area("Raw AI Output (for debugging):", value=response, height=200)
                    except Exception as e:
                        st.error(f"Failed to create .docx file. Error: {e}")
                        st.text_area("Raw AI Output (for debugging):", value=response, height=200)

    # --- TAB 2: POWERPOINT GENERATOR ---
    with tab_ppt:
        st.header("PowerPoint Generator (.pptx)")
        gem_name = "PowerPoint Generator (.pptx)"
        
        with st.expander("View/Edit Gem Prompt"):
            st.session_state.gems[gem_name] = st.text_area(
                "Prompt for PPT Generator:",
                value=st.session_state.gems.get(gem_name),
                height=300,
                key="gem_ppt"
            )

        uploaded_file = st.file_uploader("Upload your textbook (PDF)", type="pdf", key="ppt_uploader")
        
        if st.button("Generate PowerPoint"):
            if not uploaded_file:
                st.warning("Please upload a PDF file.")
            elif not st.session_state.api_key:
                st.error("API Key not set.")
            else:
                with st.spinner(f"Reading {uploaded_file.name}..."):
                    try:
                        with pdfplumber.open(uploaded_file) as pdf:
                            text_content = ""
                            for page in pdf.pages[:10]: 
                                text_content += page.extract_text() + "\n"
                        st.info("Reading first 10 pages of PDF...")
                    except Exception as e:
                        st.error(f"Failed to read PDF. Error: {e}")
                        text_content = None

                if text_content:
                    with st.spinner("Calling Gemini to generate slide content (this may take a moment)..."):
                        response = call_gemini_api(
                            st.session_state.gems[gem_name],
                            text_content,
                            st.session_state.temperature,
                            st.session_state.model_name
                        )
                    
                    if response:
                        try:
                            st.write("  > AI returned content. Parsing JSON...")
                            
                            json_start = response.find('{')
                            json_end = response.rfind('}') + 1
                            if json_start == -1 or json_end == -1:
                                raise json.JSONDecodeError("No JSON object found in response.", response, 0)
                            response_clean = response[json_start:json_end]
                            slide_data = json.loads(response_clean)
                            
                            from pptx import Presentation
                            
                            prs = Presentation()
                            
                            for slide_info in slide_data.get("slides", []):
                                slide_layout = prs.slide_layouts[1] # 1 is "Title and Content"
                                slide = prs.slides.add_slide(slide_layout)
                                
                                if slide.shapes.title:
                                    slide.shapes.title.text = slide_info.get("title", "No Title")
                                
                                content_frame = slide.placeholders[1].text_frame
                                content_frame.clear() 
                                
                                body_list = slide_info.get("body", [])
                                if body_list:
                                    p = content_frame.paragraphs[0]
                                    p.text = body_list[0]
                                    for body_item in body_list[1:]:
                                        p = content_frame.add_paragraph()
                                        p.text = body_item
                            
                            bio = io.BytesIO()
                            prs.save(bio)
                            
                            st.success("PowerPoint generated!")
                            st.download_button(
                                label="Download PowerPoint as .pptx",
                                data=bio.getvalue(),
                                file_name=f"{uploaded_file.name}_presentation.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                        
                        except json.JSONDecodeError:
                            st.error("AI returned invalid JSON. Could not create .pptx.")
                            st.text_area("Raw AI Output:", value=response, height=200)
                        except Exception as e:
                            st.error(f"Failed to create .pptx file. Error: {e}")
                            st.text_area("Raw AI Output:", value=response, height=200)

    # --- TAB 3: GIMKIT GENERATOR ---
    with tab_gimkit:
        st.header("Gimkit Generator (.csv)")
        gem_name = "Gimkit Generator (.csv)"
        
        with st.expander("View/Edit Gem Prompt"):
            st.session_state.gems[gem_name] = st.text_area(
                "Prompt for Gimkit Generator:",
                value=st.session_state.gems.get(gem_name),
                height=300,
                key="gem_gimkit"
            )
        
        topic_gimkit = st.text_input("Topic for questions:", "e.g., VCE Research Methods")
        num_gimkit = st.number_input("Number of questions:", min_value=5, max_value=50, value=15)
        
        if st.button("Generate Gimkit CSV"):
            if not topic_gimkit:
                st.warning("Please enter a topic.")
            elif not st.session_state.api_key:
                st.error("API Key not set.")
            else:
                user_prompt = f"Topic: {topic_gimkit}\nNumber of Questions: {num_gimkit}"
                with st.spinner("Generating questions..."):
                    response = call_gemini_api(
                        st.session_state.gems[gem_name],
                        user_prompt,
                        st.session_state.temperature,
                        st.session_state.model_name
                    )
                
                if response:
                    st.success("CSV content generated!")
                    
                    response_clean = re.sub(r'```csv\n(.*?)\n```', r'\1', response, flags=re.DOTALL)
                    response_clean = response_clean.strip() 
                    
                    st.download_button(
                        label="Download Gimkit CSV",
                        data=response_clean.encode('utf-8'),
                        file_name=f"{topic_gimkit.replace(' ', '_')}_gimkit.csv",
                        mime="text/csv"
                    )
                    with st.expander("Preview raw CSV data"):
                        st.text(response_clean)

    # --- TAB 4: COMMENT GENERATOR ---
    with tab_comment:
        st.header("Rubric Comment Generator")
        gem_name = "Rubric Comment Generator"
        
        with st.expander("View/Edit Gem Prompt"):
            st.session_state.gems[gem_name] = st.text_area(
                "Prompt for Comment Generator:",
                value=st.session_state.gems.get(gem_name),
                height=300,
                key="gem_comment"
            )
        
        col1, col2 = st.columns(2)
        with col1:
            questions = st.text_area("Test Questions:", "e.g., Q1: Define 'IV'.\nQ2: Explain 'extraneous variable'.", height=200)
            rubric = st.text_area("Rubric Criteria:", "e.g., 'A: Clearly defines all terms.'\n'B: Provides accurate examples.'", height=200)
        with col2:
            student_performance = st.text_area(
                "Student's Performance:", 
                "e.g., 'Correctly defined IV in Q1, but confused 'extraneous' with 'confounding' in Q2. Did not provide an example.'", 
                height=410
            )

        if st.button("Generate Comment"):
            if not questions or not rubric or not student_performance:
                st.warning("Please fill in all three fields.")
            elif not st.session_state.api_key:
                st.error("API Key not set.")
            else:
                user_prompt = f"TEST QUESTIONS:\n{questions}\n\nRUBRIC:\n{rubric}\n\nSTUDENT PERFORMANCE:\n{student_performance}"
                with st.spinner("Generating comment..."):
                    response = call_gemini_api(
                        st.session_state.gems[gem_name],
                        user_prompt,
                        st.session_state.temperature,
                        st.session_state.model_name
                    )
                
                if response:
                    st.success("Comment generated!")
                    st.text_area("Generated Comment:", value=response, height=200)

    # --- TAB 5: GEM CREATOR (MANAGE GEMS) ---
    with tab_gems:
        st.header("Gem Creator (Manage Prompts)")
        st.write("Here you can create, edit, or delete the system prompts (Gems) used by the other tabs.")

        gem_list = list(st.session_state.gems.keys())
        
        selected_gem_name = st.selectbox("Select a Gem to Edit or Delete", options=[""] + gem_list, key="gem_editor_select")

        if selected_gem_name:
            current_name = selected_gem_name
            current_prompt = st.session_state.gems.get(selected_gem_name, "")
            
            gem_name = st.text_input("Gem Name", value=current_name)
            gem_prompt = st.text_area("Gem Prompt", value=current_prompt, height=300)
            
            col1, col2 = st.columns([1, 1])
            with col1:
                if st.button("Update Gem"):
                    if current_name != gem_name: 
                        if gem_name in st.session_state.gems:
                            st.error("A Gem with that new name already exists.")
                        else:
                            del st.session_state.gems[current_name]
                            st.session_state.gems[gem_name] = gem_prompt
                            save_to_file("gems.json", st.session_state.gems)
                            st.success(f"Updated '{gem_name}'!")
                            st.rerun()
                    else:
                        st.session_state.gems[gem_name] = gem_prompt
                        save_to_file("gems.json", st.session_state.gems)
                        st.success(f"Updated '{gem_name}'!")
                        st.rerun()
            with col2:
                if st.button("Delete Gem", type="primary"):
                    if current_name == "Blank Chat Prompt":
                        st.error("Cannot delete the default 'Blank Chat Prompt'.")
                    else:
                        del st.session_state.gems[current_name]
                        save_to_file("gems.json", st.session_state.gems)
                        st.success(f"Deleted '{current_name}'!")
                        st.rerun()
            
        else:
            st.subheader("Create a New Gem")
            gem_name = st.text_input("New Gem Name", key="new_gem_name")
            gem_prompt = st.text_area("New Gem Prompt", key="new_gem_prompt", height=300)
            
            if st.button("Save New Gem"):
                if gem_name and gem_prompt:
                    if gem_name in st.session_state.gems:
                        st.error("A Gem with this name already exists.")
                    else:
                        st.session_state.gems[gem_name] = gem_prompt
                        save_to_file("gems.json", st.session_state.gems)
                        st.success(f"Saved new Gem: '{gem_name}'!")
                        st.rerun()
                else:
                    st.warning("Please provide both a name and a prompt.")

# --- (6) THIS BLOCK RUNS THE SCRIPT ---
if __name__ == "__main__":
    main()
